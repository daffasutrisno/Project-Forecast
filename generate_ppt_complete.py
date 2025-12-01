"""
Generate Complete PowerPoint Presentation
Mencakup SEMUA hasil forecast dan analysis
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import pandas as pd

# Directories
MAIN_DIR = Path("forecast_results/01_main")
REGIONAL_DIR = Path("forecast_results/02_regional")
PROVINSI_DIR = Path("forecast_results/03_provinsi")
KABUPATEN_DIR = Path("forecast_results/04_kabupaten")
ANALYSIS_DIR = Path("forecast_results/05_analysis")

def create_slide_with_image(prs, image_path, title_text):
    """Create a slide with an image and title"""
    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(0.6)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title_text
    
    # Format title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    
    # Add image
    if Path(image_path).exists():
        img_left = Inches(0.5)
        img_top = Inches(1)
        img_width = Inches(9)
        slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width)
    else:
        print(f"  ⚠️  Image tidak ditemukan: {image_path}")
    
    return slide

def get_all_kabupaten_images():
    """Get all kabupaten forecast images sorted alphabetically"""
    kabupatens = []
    
    if KABUPATEN_DIR.exists():
        for img_file in sorted(KABUPATEN_DIR.glob("*_forecast.png")):
            # Extract kabupaten name from filename
            kabupaten_name = img_file.stem.replace("_forecast", "").replace("_", " ").upper()
            kabupatens.append((img_file, kabupaten_name))
    
    return kabupatens

def create_presentation():
    """Create Complete PowerPoint presentation"""
    
    print("="*80)
    print("  GENERATE COMPLETE POWERPOINT PRESENTATION")
    print("="*80)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide_count = 0
    
    # ========================================================================
    # SECTION 1: MAIN FORECASTS
    # ========================================================================
    print("\n📊 SECTION 1: Main Forecasts")
    print("-" * 80)
    
    # Slide: Traffic Forecast Lengkap (Combined provinces)
    img_path = PROVINSI_DIR / "04_combined_province_comparison.png"
    if img_path.exists():
        create_slide_with_image(prs, img_path, "Traffic Forecast Lengkap - Perbandingan Provinsi")
        slide_count += 1
        print(f"  ✓ Slide {slide_count}: Traffic Forecast Lengkap")
    
    # Slide: Main Forecast Overview
    img_path = MAIN_DIR / "00_main_forecast_overview.png"
    if img_path.exists():
        create_slide_with_image(prs, img_path, "Main Forecast Overview - Total Traffic")
        slide_count += 1
        print(f"  ✓ Slide {slide_count}: Main Forecast Overview")
    
    # ========================================================================
    # SECTION 2: REGIONAL FORECASTS
    # ========================================================================
    print("\n📊 SECTION 2: Regional Forecasts (3 regions)")
    print("-" * 80)
    
    regions = [
        ("bali_nusra.csv", "Forecast Regional: Bali Nusra"),
        ("central_java.csv", "Forecast Regional: Central Java"),
        ("east_java.csv", "Forecast Regional: East Java")
    ]
    
    for csv_file, title in regions:
        img_path = REGIONAL_DIR / csv_file.replace(".csv", "_forecast.png")
        if img_path.exists():
            create_slide_with_image(prs, img_path, title)
            slide_count += 1
            print(f"  ✓ Slide {slide_count}: {title}")
    
    # ========================================================================
    # SECTION 3: PROVINSI FORECASTS
    # ========================================================================
    print("\n📊 SECTION 3: Provinsi Forecasts (6 provinces)")
    print("-" * 80)
    
    provinces = [
        ("bali.csv", "Forecast Provinsi: Bali"),
        ("daerah_istimewa_yogyakarta.csv", "Forecast Provinsi: Daerah Istimewa Yogyakarta"),
        ("jawa_tengah.csv", "Forecast Provinsi: Jawa Tengah"),
        ("jawa_timur.csv", "Forecast Provinsi: Jawa Timur"),
        ("nusa_tenggara_barat.csv", "Forecast Provinsi: Nusa Tenggara Barat"),
        ("nusa_tenggara_timur.csv", "Forecast Provinsi: Nusa Tenggara Timur")
    ]
    
    for csv_file, title in provinces:
        img_path = PROVINSI_DIR / csv_file.replace(".csv", "_forecast.png")
        if img_path.exists():
            create_slide_with_image(prs, img_path, title)
            slide_count += 1
            print(f"  ✓ Slide {slide_count}: {title}")
    
    # ========================================================================
    # SECTION 4: ALL KABUPATEN FORECASTS (119 kabupatens)
    # ========================================================================
    print("\n📊 SECTION 4: Kabupaten Forecasts (119 kabupatens)")
    print("-" * 80)
    
    kabupatens = get_all_kabupaten_images()
    print(f"  📍 Total kabupaten ditemukan: {len(kabupatens)}")
    
    for img_path, kabupaten_name in kabupatens:
        create_slide_with_image(prs, img_path, f"Forecast Kabupaten: {kabupaten_name}")
        slide_count += 1
        if slide_count % 10 == 0:
            print(f"  ✓ Progress: {slide_count} slides (kabupaten ke-{slide_count - 11})")
    
    print(f"  ✓ Semua {len(kabupatens)} kabupaten berhasil ditambahkan")
    
    # ========================================================================
    # SECTION 5: TOP 10 ANALYSIS - ABSOLUTE CHANGE
    # ========================================================================
    print("\n📊 SECTION 5: Top 10 Analysis - Absolute Change")
    print("-" * 80)
    
    # Summary slide
    img_path = ANALYSIS_DIR / "top10_kabupaten_absolute_summary.png"
    if img_path.exists():
        create_slide_with_image(prs, img_path, "Top 10 Kabupaten - Absolute Traffic Increase")
        slide_count += 1
        print(f"  ✓ Slide {slide_count}: Top 10 Summary (Absolute)")
    
    # Individual slides for top 10 absolute
    excel_file = ANALYSIS_DIR / "top10_kabupaten_absolute.xlsx"
    if excel_file.exists():
        df = pd.read_excel(excel_file, sheet_name='Top 10')
        for idx, row in df.iterrows():
            kabupaten_name = row['Kabupaten'].lower().replace(' ', '_')
            img_path = KABUPATEN_DIR / f"{kabupaten_name}_forecast.png"
            
            if img_path.exists():
                change_tb = row['Change_TB']
                growth_pct = row['Growth_%']
                title = f"Top #{idx+1} (Absolute): {row['Kabupaten']} (+{change_tb:.2f} TB, +{growth_pct:.2f}%)"
                create_slide_with_image(prs, img_path, title)
                slide_count += 1
        
        print(f"  ✓ 10 kabupaten individual (absolute) ditambahkan")
    
    # ========================================================================
    # SECTION 6: TOP 10 ANALYSIS - PERCENTAGE GROWTH
    # ========================================================================
    print("\n📊 SECTION 6: Top 10 Analysis - Percentage Growth")
    print("-" * 80)
    
    # Summary slide
    img_path = ANALYSIS_DIR / "top10_kabupaten_percentage_summary.png"
    if img_path.exists():
        create_slide_with_image(prs, img_path, "Top 10 Kabupaten - Highest Growth Percentage")
        slide_count += 1
        print(f"  ✓ Slide {slide_count}: Top 10 Summary (Percentage)")
    
    # Individual slides for top 10 percentage
    excel_file = ANALYSIS_DIR / "top10_kabupaten_individual_forecast.xlsx"
    if excel_file.exists():
        df = pd.read_excel(excel_file, sheet_name='Top 10')
        for idx, row in df.iterrows():
            kabupaten_name = row['Kabupaten'].lower().replace(' ', '_')
            img_path = KABUPATEN_DIR / f"{kabupaten_name}_forecast.png"
            
            if img_path.exists():
                growth_rate = row['Growth_Rate']
                title = f"Top #{idx+1} (Percentage): {row['Kabupaten']} (+{growth_rate:.2f}%)"
                create_slide_with_image(prs, img_path, title)
                slide_count += 1
        
        print(f"  ✓ 10 kabupaten individual (percentage) ditambahkan")
    
    # ========================================================================
    # SAVE PRESENTATION
    # ========================================================================
    output_file = "presentation_traffic_forecast_COMPLETE.pptx"
    prs.save(output_file)
    
    print("\n" + "="*80)
    print("✅ PRESENTASI LENGKAP BERHASIL DIBUAT!")
    print("="*80)
    print(f"\n📁 File: {output_file}")
    print(f"📊 Total Slides: {slide_count}")
    print(f"\n📋 Breakdown:")
    print(f"  • Main Forecasts: 2 slides")
    print(f"  • Regional Forecasts: 3 slides")
    print(f"  • Provinsi Forecasts: 6 slides")
    print(f"  • Kabupaten Forecasts: {len(kabupatens)} slides")
    print(f"  • Top 10 Absolute: 11 slides (1 summary + 10 individual)")
    print(f"  • Top 10 Percentage: 11 slides (1 summary + 10 individual)")
    print("="*80)

def main():
    """Main function"""
    try:
        create_presentation()
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
