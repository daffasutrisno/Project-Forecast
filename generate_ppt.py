"""
GENERATE POWERPOINT PRESENTATION
=================================
Script untuk membuat presentasi PowerPoint dari hasil forecast

Output: presentation_traffic_forecast.pptx

Author: Traffic Analysis Team
Date: November 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import pandas as pd

def create_slide_with_image(prs, image_path, title):
    """Create a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    # Format title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = 'Calibri'
    
    # Add image
    left = Inches(0.5)
    top = Inches(1.1)
    width = Inches(9)
    
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), left, top, width=width)
        return True
    else:
        print(f"  ⚠️  File tidak ditemukan: {image_path}")
        return False

def get_regional_images():
    """Get all regional forecast images"""
    regional_dir = Path("forecast_results/02_regional")
    images = []
    
    regions = ['bali_nusra', 'central_java', 'east_java']
    for region in regions:
        img_path = regional_dir / f"{region}_forecast.png"
        if img_path.exists():
            images.append((img_path, region.replace('_', ' ').title()))
    
    return images

def get_province_images():
    """Get all province forecast images"""
    province_dir = Path("forecast_results/03_provinsi")
    images = []
    
    provinces = [
        'bali',
        'daerah_istimewa_yogyakarta',
        'jawa_tengah',
        'jawa_timur',
        'nusa_tenggara_barat',
        'nusa_tenggara_timur'
    ]
    
    for province in provinces:
        img_path = province_dir / f"{province}_forecast.png"
        if img_path.exists():
            title = province.replace('_', ' ').title()
            images.append((img_path, title))
    
    return images

def get_top10_absolute_kabupatens():
    """Get top 10 kabupaten by absolute change"""
    excel_file = Path("forecast_results/05_analysis/top10_kabupaten_by_absolute_change.xlsx")
    
    if not excel_file.exists():
        print("  ⚠️  File top10_kabupaten_by_absolute_change.xlsx tidak ditemukan!")
        return []
    
    df = pd.read_excel(excel_file, sheet_name='Top 10 Absolute')
    
    kabupatens = []
    for idx, row in df.iterrows():
        kabupaten_name = row['Kabupaten'].lower().replace(' ', '_')
        img_path = Path(f"forecast_results/04_kabupaten/{kabupaten_name}_forecast.png")
        
        if img_path.exists():
            title = f"#{idx+1}: {row['Kabupaten']} ({row['Region']})"
            kabupatens.append((img_path, title))
    
    return kabupatens

def get_top10_percentage_kabupatens():
    """Get top 10 kabupaten by percentage growth"""
    excel_file = Path("forecast_results/05_analysis/top10_kabupaten_individual_forecast.xlsx")
    
    if not excel_file.exists():
        print("  ⚠️  File top10_kabupaten_individual_forecast.xlsx tidak ditemukan!")
        return []
    
    df = pd.read_excel(excel_file, sheet_name='Top 10')
    
    kabupatens = []
    for idx, row in df.iterrows():
        kabupaten_name = row['Kabupaten'].lower().replace(' ', '_')
        img_path = Path(f"forecast_results/04_kabupaten/{kabupaten_name}_forecast.png")
        
        if img_path.exists():
            growth_rate = row['Growth_Rate']  # Column name is Growth_Rate, not Growth %
            title = f"#{idx+1}: {row['Kabupaten']} (+{growth_rate:.2f}%)"
            kabupatens.append((img_path, title))
    
    return kabupatens

def create_presentation():
    """Create PowerPoint presentation"""
    
    print("="*80)
    print("  GENERATE POWERPOINT PRESENTATION")
    print("="*80)
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide_count = 0
    
    # 1. Traffic Forecast Lengkap
    print("\n📊 Slide 1: Traffic Forecast Lengkap")
    img_path = "forecast_results/01_main/03_traffic_forecast_lengkap.png"
    if create_slide_with_image(prs, img_path, "Traffic Forecast Lengkap: Historical + Forecast"):
        slide_count += 1
        print(f"  ✓ Slide {slide_count} ditambahkan")
    
    # 2. Main Forecast Overview
    print("\n📊 Slide 2: Main Forecast Overview")
    img_path = "forecast_results/01_main/00_main_forecast_overview.png"
    if create_slide_with_image(prs, img_path, "Main Forecast Overview: Total Java"):
        slide_count += 1
        print(f"  ✓ Slide {slide_count} ditambahkan")
    
    # 3. Regional Forecasts
    print("\n📊 Slides 3-5: Forecast Regional")
    regional_images = get_regional_images()
    for img_path, title in regional_images:
        if create_slide_with_image(prs, img_path, f"Forecast Regional: {title}"):
            slide_count += 1
            print(f"  ✓ Slide {slide_count} ditambahkan: {title}")
    
    # 4. Province Forecasts
    print("\n📊 Slides 6-11: Forecast Provinsi")
    province_images = get_province_images()
    for img_path, title in province_images:
        if create_slide_with_image(prs, img_path, f"Forecast Provinsi: {title}"):
            slide_count += 1
            print(f"  ✓ Slide {slide_count} ditambahkan: {title}")
    
    # 5. Top 10 Kabupaten Absolute - Summary
    print("\n📊 Slide: Top 10 Kabupaten by Absolute Change")
    img_path = "forecast_results/05_analysis/top10_kabupaten_by_absolute_change.png"
    if create_slide_with_image(prs, img_path, "Top 10 Kabupaten: Peningkatan Absolut (TB)"):
        slide_count += 1
        print(f"  ✓ Slide {slide_count} ditambahkan")
    
    # 6. Top 10 Kabupaten Absolute - Individual
    print("\n📊 Slides: Top 10 Kabupaten Individual (Absolute)")
    top10_absolute = get_top10_absolute_kabupatens()
    for img_path, title in top10_absolute:
        if create_slide_with_image(prs, img_path, f"Top 10 Absolute: {title}"):
            slide_count += 1
            print(f"  ✓ Slide {slide_count} ditambahkan")
    
    # 7. Top 10 Kabupaten Percentage - Summary
    print("\n📊 Slide: Top 10 Kabupaten by Percentage Growth")
    img_path = "forecast_results/05_analysis/top10_kabupaten_individual_forecast.png"
    if create_slide_with_image(prs, img_path, "Top 10 Kabupaten: Growth Percentage (%)"):
        slide_count += 1
        print(f"  ✓ Slide {slide_count} ditambahkan")
    
    # 8. Top 10 Kabupaten Percentage - Individual
    print("\n📊 Slides: Top 10 Kabupaten Individual (Percentage)")
    top10_percentage = get_top10_percentage_kabupatens()
    for img_path, title in top10_percentage:
        if create_slide_with_image(prs, img_path, f"Top 10 Percentage: {title}"):
            slide_count += 1
            print(f"  ✓ Slide {slide_count} ditambahkan")
    
    # Save presentation
    output_file = "presentation_traffic_forecast.pptx"
    prs.save(output_file)
    
    print("\n" + "="*80)
    print("✅ PRESENTASI BERHASIL DIBUAT!")
    print("="*80)
    print(f"\n📁 File: {output_file}")
    print(f"📊 Total Slides: {slide_count}")
    print("\nStruktur Presentasi:")
    print("  1. Traffic Forecast Lengkap")
    print("  2. Main Forecast Overview")
    print(f"  3-5. Regional Forecasts ({len(regional_images)} slides)")
    print(f"  6-11. Province Forecasts ({len(province_images)} slides)")
    print("  12. Top 10 Kabupaten by Absolute (Summary)")
    print(f"  13-22. Top 10 Kabupaten by Absolute (Individual - {len(top10_absolute)} slides)")
    print("  23. Top 10 Kabupaten by Percentage (Summary)")
    print(f"  24-33. Top 10 Kabupaten by Percentage (Individual - {len(top10_percentage)} slides)")
    print("\n" + "="*80)

def main():
    """Main function"""
    try:
        create_presentation()
    except Exception as e:
        print(f"\n❌ Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
